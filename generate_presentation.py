"""
Script to generate the powerpoint presentation
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
import gorilla

import content_placeholder_patch


# first, monkey patch using gorilla as per https://github.com/scanny/python-pptx/issues/333#issuecomment-427055526
patches = gorilla.find_patches([content_placeholder_patch])
for patch in patches:
    gorilla.apply(patch)


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "from pptx import Presentation"
subtitle.text = "Building PowerPoint decks with Python"

# Layouts:
#
# 0.   Title Slide,
# 1.   Title and Content,
# 2.   Section Header,
# 3.   Two Content,
# 4.   Comparison,
# 5.   Title Only,
# 6.   Blank,
# 7.   Content with Caption,
# 8.   Picture with Caption,
# 9.   Title and Vertical Text,
# 10.  Vertical Title and Text

title_content_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title
title.text = 'Have you ever been in this position?'
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = 'Every week you write a PowerPoint presentation'
p = tf.add_paragraph()
p.text = "It's not long; you use it to give your project status report"
p = tf.add_paragraph()
p.text = "You're on a long project, so the slides don't change much week to week"
p = tf.add_paragraph()
p.text = "Even though all the text comes from a standard template..."


section_header_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(section_header_layout)
title = slide.shapes.title
title.text = "You're a developer. It doesn't have to be this way"


two_content_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = 'Python to the rescue (again)'
left_body = slide.shapes.placeholders[1]
tf = left_body.text_frame
tf.text = 'The pptx library lets you make a PowerPoint presentation in Python'
p = tf.add_paragraph()
run = p.add_run()
run.text = "(though it does require a "
run = p.add_run()
font = run.font
font.italic = True
run.text = "tiny"
run = p.add_run()
run.text = " monkey-patch at the moment)"

right_body = slide.shapes.placeholders[2]
phf = right_body.placeholder_format
right_body.insert_picture('this_is_fine.jpeg')


slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = 'Great! So how does it work?'

left_body = slide.shapes.placeholders[1]
tf = left_body.text_frame
tf.text = "It's pretty good, if a little verbose"
p = tf.add_paragraph()
p.text = 'The first step is to create the presentation object'
p = tf.add_paragraph()
p.text = "You're looking at the default presentation. If you pass a filename you can work of any template"


right_body = slide.shapes.placeholders[2]
tf = right_body.text_frame
p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'from pptx import Presentation'
p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'prs = Presentation()'


slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "From there it's all adding content"

left_body = slide.shapes.placeholders[1]
tf = left_body.text_frame
p = tf.add_paragraph()
p.text = "You can add slides based on the master slide templates"
p = tf.add_paragraph()
p.text = "The default prder of these masters is always the same"
p = tf.add_paragraph()
p.text = "And then you can add content"

right_body = slide.shapes.placeholders[2]
tf = right_body.text_frame
p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'prs = Presentation()'

p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'title_slide_layout = prs.slide_layouts[0]'

p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'slide = prs.slides.add_slide(title_slide_layout)'

p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'title = slide.shapes.title'

p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'subtitle = slide.placeholders[1]'

p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'title.text = "from pptx import Presentation"'

p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = 'subtitle.text = "Building PowerPoint decks with Python"'

slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title
title.text = "What's supported?"

body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = 'Round-trip any Open XML presentation (.pptx file) including all its elements'
p = tf.add_paragraph()
p.text = 'Add slides'
p = tf.add_paragraph()
p.text = 'Populate text placeholders, for example to create a bullet slide'
p = tf.add_paragraph()
p.text = 'Add image to slide at arbitrary position and size'
p = tf.add_paragraph()
p.text = 'Add textbox to a slide; manipulate text font size and bold'
p = tf.add_paragraph()
p.text = 'Add table to a slide'
p = tf.add_paragraph()
p.text = 'Add auto shapes (e.g. polygons, flowchart shapes, etc.) to a slide'
p = tf.add_paragraph()
p.text = 'Add and manipulate column, bar, line, and pie charts'
p = tf.add_paragraph()
p.text = 'Access and change core document properties such as title and subject'


title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding an AutoShape'

left = Inches(0.93)  # 0.93" centers this overall set of shapes
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = 'Step 1'

left = left + width - Inches(0.4)
width = Inches(2.0)  # chevrons need more width for visual balance

for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = 'Step %d' % n
    left = left + width - Inches(0.4)

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "Finally saving the presentation"

left_body = slide.shapes.placeholders[1]
tf = left_body.text_frame
p = tf.add_paragraph()
p.text = "Once you've finished adding everything, all that's left is to save"
p = tf.add_paragraph()
p.text = '(And yes, this presentation was all written in Python)'
p = tf.add_paragraph()
p.level = 1
p.text = '245 lines of code'
p = tf.add_paragraph()
run = p.add_run()
run.text = 'https://github.com/om-henners/mpug_march_2019'
run.hyperlink.address = 'https://github.com/om-henners/mpug_march_2019'

right_body = slide.shapes.placeholders[2]
tf = right_body.text_frame
p = tf.add_paragraph()
p.font.name = 'Source Code Pro for Powerline'
p.text = "prs.save('mpug_march_2019.pptx')"


prs.save('mpug_march_2019.pptx')
