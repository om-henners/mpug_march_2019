"""
Monkey patch for python-pptx that allows images to be placed into content
placeholders.
"""

import gorilla
import pptx
from pptx.oxml.shapes.picture import CT_Picture
from pptx.shapes.placeholder import _BaseSlidePlaceholder, PlaceholderPicture

settings = gorilla.Settings(allow_hit=True, store_hit=True)
@gorilla.patches(pptx.shapes.placeholder.SlidePlaceholder, settings=settings)
class SlidePlaceholder(_BaseSlidePlaceholder):
    """
    Placeholder shape on a slide. Inherits shape properties from its
    corresponding slide layout placeholder.
    """
    def insert_picture(self, image_file, crop=True):
        """
        Return a |PlaceholderPicture| object depicting the image in
        *image_file*, which may be either a path (string) or a file-like
        object. The image is cropped to fill the entire space of the
        placeholder. A |PlaceholderPicture| object has all the properties and
        methods of a |Picture| shape except that the value of its
        :attr:`~._BaseSlidePlaceholder.shape_type` property is
        `MSO_SHAPE_TYPE.PLACEHOLDER` instead of `MSO_SHAPE_TYPE.PICTURE`.
        """
        pic = self.new_placeholder_pic(image_file, crop) # pass new parameter "method"
        self._replace_placeholder_with(pic)
        return PlaceholderPicture(pic, self._parent)

    def new_placeholder_pic(self, image_file, crop=True):
        """
        Return a new `p:pic` element depicting the image in *image_file*,
        suitable for use as a placeholder. In particular this means not
        having an `a:xfrm` element, allowing its extents to be inherited from
        its layout placeholder.
        """
        rId, desc, image_size = self.get_or_add_image(image_file)
        id_, name = self.shape_id, self.name

        # Cropping the image, as in the original file
        if crop:
            pic = CT_Picture.new_ph_pic(id_, name, desc, rId)
            pic.crop_to_fit(image_size, (self.width, self.height))

        # Adjusting image to placeholder size and replace placeholder.
        else:
            ph_w, ph_h = self.width, self.height
            aspectPh = ph_w / ph_h

            img_w, img_h = image_size
            aspectImg = img_w / img_h

            if aspectPh > aspectImg:
                w = int(ph_h * aspectImg)
                h = ph_h # keep the height
            else:
                w = ph_w # keep the width
                h = int(ph_w / aspectImg)

            top = self.top + (ph_h - h) / 2
            left = self.left + (ph_w - w) / 2

            pic = CT_Picture.new_pic(id_, name, desc, rId, self.left + (ph_w - w) / 2, self.top, w, h)
        return pic

    def get_or_add_image(self, image_file):
        """
        Return an (rId, description, image_size) 3-tuple identifying the
        related image part containing *image_file* and describing the image.
        """
        image_part, rId = self.part.get_or_add_image_part(image_file)
        desc, image_size = image_part.desc, image_part._px_size
        return rId, desc, image_size
